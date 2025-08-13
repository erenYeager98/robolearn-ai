import os
import json
import fitz  # PyMuPDF
import docx
import pptx
import yake
import logging

# Configure basic logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

# --- CONFIGURATION ---
# IMPORTANT: Specify the path to the folder you want to scan.
# Example for Windows: 'C:\\Users\\YourUser\\Documents\\MyFiles'
# Example for macOS/Linux: '/home/youruser/documents/myfiles'
SOURCE_FOLDER = 'C:\\Users\\UTASI\\OneDrive\\Desktop\\PI'

# Output files and folders
OUTPUT_JSON_FILE = 'extracted_content.json'
KEYWORDS_FOLDER = 'local_keywords'
KEYWORDS_FILE = os.path.join(KEYWORDS_FOLDER, 'keywords.txt')
# --- END CONFIGURATION ---


def extract_text_from_pdf(file_path: str) -> str:
    """Extracts text from a PDF file."""
    try:
        with fitz.open(file_path) as doc:
            text = "".join(page.get_text() for page in doc)
        logging.info(f"Successfully extracted text from PDF: {os.path.basename(file_path)}")
        return text
    except Exception as e:
        logging.error(f"Could not read PDF {os.path.basename(file_path)}: {e}")
        return ""

def extract_text_from_docx(file_path: str) -> str:
    """Extracts text from a DOCX file."""
    try:
        doc = docx.Document(file_path)
        text = "\n".join(para.text for para in doc.paragraphs)
        logging.info(f"Successfully extracted text from DOCX: {os.path.basename(file_path)}")
        return text
    except Exception as e:
        logging.error(f"Could not read DOCX {os.path.basename(file_path)}: {e}")
        return ""

def extract_text_from_pptx(file_path: str) -> str:
    """Extracts text from a PPTX file."""
    try:
        prs = pptx.Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        logging.info(f"Successfully extracted text from PPTX: {os.path.basename(file_path)}")
        return text
    except Exception as e:
        logging.error(f"Could not read PPTX {os.path.basename(file_path)}: {e}")
        return ""

def process_folder(folder_path: str) -> dict:
    """
    Recursively scans a folder and extracts text from supported files.
    """
    if not os.path.isdir(folder_path):
        logging.error(f"Source folder not found: {folder_path}")
        return {}
        
    all_docs_content = {}
    supported_extensions = {
        '.pdf': extract_text_from_pdf,
        '.docx': extract_text_from_docx,
        '.pptx': extract_text_from_pptx,
    }

    logging.info(f"Starting to scan folder: {folder_path}")
    for root, _, files in os.walk(folder_path):
        for file in files:
            file_path = os.path.join(root, file)
            _, extension = os.path.splitext(file)

            if extension in supported_extensions:
                extractor_func = supported_extensions[extension]
                content = extractor_func(file_path)
                if content:
                    all_docs_content[file_path] = content
            else:
                logging.warning(f"Skipping unsupported file type: {file}")
                
    return all_docs_content

def generate_keywords(full_text: str, num_keywords: int = 50) -> list:
    """
    Generates a list of top keywords from the given text.
    """
    if not full_text.strip():
        logging.warning("No text content found to generate keywords from.")
        return []
    
    logging.info("Generating keywords from the extracted content...")
    # Using YAKE! for keyword extraction
    # n = max n-gram size, dedupLim = deduplication threshold, top = number of keywords
    kw_extractor = yake.KeywordExtractor(lan="en", n=3, dedupLim=0.9, top=num_keywords, features=None)
    keywords = kw_extractor.extract_keywords(full_text)
    
    # Return just the keyword strings
    return [kw for kw, score in keywords]

def main():
    """Main function to run the script."""
    # 1. Extract text from all documents in the folder
    extracted_data = process_folder(SOURCE_FOLDER)

    if not extracted_data:
        logging.info("No text was extracted. Exiting.")
        return

    # 2. Store the extracted content in a JSON file
    logging.info(f"Saving all extracted text to {OUTPUT_JSON_FILE}...")
    try:
        with open(OUTPUT_JSON_FILE, 'w', encoding='utf-8') as f:
            json.dump(extracted_data, f, indent=4, ensure_ascii=False)
        logging.info("Successfully saved content to JSON.")
    except Exception as e:
        logging.error(f"Failed to write to JSON file: {e}")
        return

    # 3. Combine all text for keyword analysis
    combined_text = " ".join(extracted_data.values())

    # 4. Generate keywords
    keywords_list = generate_keywords(combined_text)

    if not keywords_list:
        logging.info("No keywords were generated. Exiting keyword saving process.")
        return

    # 5. Store keywords in a text file under the 'local_keywords' folder
    logging.info(f"Saving keywords to {KEYWORDS_FILE}...")
    try:
        os.makedirs(KEYWORDS_FOLDER, exist_ok=True)
        with open(KEYWORDS_FILE, 'w', encoding='utf-8') as f:
            f.write("\n".join(keywords_list))
        logging.info(f"Successfully saved {len(keywords_list)} keywords.")
    except Exception as e:
        logging.error(f"Failed to write keywords to file: {e}")

    print("\n✅ Process completed successfully!")


if __name__ == "__main__":
    main()